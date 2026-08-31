#!/usr/bin/env python3
"""
generate_deck.py - Generates Project SUTRA PowerPoint Presentation (.pptx)
Compliant with PPT Artisan Skills, Stitch Tactical Monolith & Hallmark Anti-Slop Discipline:
- Slide 01: Title Slide (Monumental headline, 3 Grounded Tabs, Team Subsystem Roster)
- Slide 02: The Problem (4-Column Subsystem Failure Voids Grid: A, B, C, D)
- Slide 03: The Solution (4-Column Subsystem Moats Architecture Grid: A, B, C, D)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Web Aura & Tactical Monolith Palette
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)        # Pure White (#FFFFFF)
COLOR_BLACK = RGBColor(0x09, 0x09, 0x0B)     # Tactical Ink Black (#09090B)
COLOR_NAVY = RGBColor(0x0A, 0x16, 0x28)      # Deep Navy Accent (#0A1628)
COLOR_GRAY = RGBColor(0x37, 0x41, 0x51)      # Dark Gray Text (#374151)
COLOR_SLATE = RGBColor(0x47, 0x55, 0x69)     # Slate Neutral (#475569)
COLOR_MUTED = RGBColor(0x64, 0x74, 0x8B)     # Muted Text (#64748B)
COLOR_DIM = RGBColor(0x94, 0xA3, 0xB8)       # Dim Label (#94A3B8)
COLOR_BORDER = RGBColor(0xE4, 0xE4, 0xE7)    # Hairline Border (#E4E4E7)
COLOR_CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)   # Slate Surface (#F8FAFC)
COLOR_WATERMARK = RGBColor(0xF1, 0xF5, 0xF9) # Faint Watermark Number (#F1F5F9)
COLOR_EMERALD = RGBColor(0x04, 0x78, 0x57)   # Emerald Status Pill
COLOR_LIGHT_EMERALD = RGBColor(0x34, 0xD3, 0x99) # Emerald Green Text (#34D399)
COLOR_RED = RGBColor(0xE1, 0x1D, 0x48)       # Crimson Warning Pill

FONT_HEADING = "Space Grotesk"
FONT_SERIF = "Georgia"
FONT_BODY = "Inter"
FONT_MONO = "JetBrains Mono"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_name=FONT_BODY,
             font_size=14, color=COLOR_BLACK, bold=False, italic=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if p.runs:
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return tx


def build_slide_01_title(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 1. Background Grid Pattern
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_BG)
    grid_spacing = 0.5
    for x_step in range(1, 26):
        x = x_step * grid_spacing
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(0), Inches(0.008), SLIDE_HEIGHT, COLOR_BORDER)
    for y_step in range(1, 15):
        y = y_step * grid_spacing
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(y), SLIDE_WIDTH, Inches(0.008), COLOR_BORDER)

    # 2. Top Header
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.55), Inches(0.35), Inches(0.35), COLOR_BLACK)
    add_text(slide, Inches(1.25), Inches(0.62), Inches(6.0), Inches(0.3),
             "TEAM OFFGRID  /  DEFENSE & DISASTER ROBOTICS", FONT_HEADING, 9.5, COLOR_BLACK, bold=True)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.58), Inches(1.8), Inches(0.3), COLOR_CARD_BG, COLOR_BORDER)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(9.95), Inches(0.68), Inches(0.08), Inches(0.08), COLOR_EMERALD)
    add_text(slide, Inches(10.1), Inches(0.64), Inches(1.5), Inches(0.2), "AUTONOMOUS SWARM", FONT_HEADING, 7.5, COLOR_SLATE, bold=True)
    add_text(slide, Inches(11.7), Inches(0.62), Inches(1.5), Inches(0.3),
             "AUG 2026 • REV 1.0", FONT_MONO, 8.5, COLOR_MUTED, align=PP_ALIGN.RIGHT)

    # 3. Center Hero
    add_text(slide, Inches(0.8), Inches(1.55), Inches(11.0), Inches(0.4),
             "“When GPS fails and RF links jam, SUTRA geonavigates and locates survivors in real-time.”",
             FONT_SERIF, 18, COLOR_BLACK, italic=True)

    bracket_len = Inches(0.35)
    bracket_thick = Inches(0.02)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.05), bracket_len, bracket_thick, COLOR_BLACK)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.05), bracket_thick, bracket_len, COLOR_BLACK)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(3.75), bracket_len, bracket_thick, COLOR_BLACK)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.33), Inches(3.42), bracket_thick, bracket_len, COLOR_BLACK)

    add_text(slide, Inches(0.8), Inches(2.05), Inches(8.0), Inches(0.8), "PROJECT", FONT_HEADING, 52, COLOR_BLACK, bold=True)
    add_text(slide, Inches(0.8), Inches(2.8), Inches(8.0), Inches(0.8), "SUTRA.", FONT_HEADING, 52, COLOR_SLATE, italic=True)

    add_text(slide, Inches(0.8), Inches(3.68), Inches(10.0), Inches(0.45),
             "Swarm Unified Tactical Reconnaissance Architecture — decentralized multi-UAV flight, Deep JSCC neural zero-cliff video, and 3.59cm terrain-corrected DEM survivor geolocation.",
             FONT_BODY, 12.0, COLOR_SLATE)

    # Grounded Tabs
    grounded_tabs = [
        ("01", "DECENTRALIZED GNC", "50Hz", "Offboard Rate", "PX4 MicroXRCE + VIO EKF2", COLOR_BLACK),
        ("02", "SEMANTIC WIRELESS", "-5dB", "Jam Resilience", "Deep JSCC Zero-Cliff Video", COLOR_NAVY),
        ("03", "TERRAIN GEOLOCATION", "3.59cm", "Geo-Accuracy", "WGS84 3D DEM Raycasting", COLOR_EMERALD),
    ]
    card_w = Inches(2.7)
    card_h = Inches(1.15)
    card_gap = Inches(0.2)
    for i, (num, tag, val, label, sub, tag_color) in enumerate(grounded_tabs):
        c_left = Inches(0.8) + i * (card_w + card_gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(4.3), card_w, card_h, COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, c_left + Inches(1.8), Inches(4.25), Inches(0.8), Inches(0.6), num, FONT_MONO, 28, COLOR_WATERMARK, bold=True, align=PP_ALIGN.RIGHT)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, c_left + Inches(0.12), Inches(4.4), Inches(1.3), Inches(0.2), tag_color)
        add_text(slide, c_left + Inches(0.12), Inches(4.42), Inches(1.3), Inches(0.18), tag, FONT_HEADING, 6.2, COLOR_BG, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, c_left + Inches(0.12), Inches(4.65), Inches(1.4), Inches(0.4), val, FONT_HEADING, 19, COLOR_BLACK, bold=True)
        add_text(slide, c_left + Inches(1.2), Inches(4.72), Inches(1.3), Inches(0.3), label, FONT_HEADING, 9, COLOR_SLATE, bold=True)
        add_text(slide, c_left + Inches(0.12), Inches(5.1), Inches(2.45), Inches(0.25), sub, FONT_BODY, 7.5, COLOR_MUTED)

    # 4. Bottom Roster Bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.65), Inches(11.733), Inches(0.015), COLOR_BORDER)
    add_text(slide, Inches(0.8), Inches(5.75), Inches(6.0), Inches(0.2), "CORE ARCHITECTURE TEAM (OFFGRID)", FONT_HEADING, 8.0, COLOR_DIM, bold=True)
    add_text(slide, Inches(7.5), Inches(5.75), Inches(5.033), Inches(0.2), "RESEARCH-BACKED & EMPIRICALLY VALIDATED", FONT_HEADING, 8.0, COLOR_SLATE, bold=True, align=PP_ALIGN.RIGHT)

    members = [
        ("Nikhil", "Tech Lead · Subsys A & B", "GNC & JSCC", True),
        ("Vedanth Sai Ram", "Lead · Subsystem C", "AI PERCEPTION", False),
        ("Siva Kesava", "Lead · Subsystem D", "3D GIS GCS", False),
        ("Harika", "Lead · Subsystem E", "VERIFICATION QA", False),
        ("Rohith Kumar", "Lead · Subsystem F", "NDMA CONOPS", False),
    ]
    col_w = Inches(2.2)
    gap = Inches(0.18)
    for i, (name, role, badge, is_lead) in enumerate(members):
        left_pos = Inches(0.8) + i * (col_w + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(6.0), col_w, Inches(1.15), COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, left_pos + Inches(0.12), Inches(6.08), col_w - Inches(0.24), Inches(0.25), name, FONT_HEADING, 10.5, COLOR_BLACK, bold=True)
        add_text(slide, left_pos + Inches(0.12), Inches(6.32), col_w - Inches(0.24), Inches(0.22), role, FONT_BODY, 8.0, COLOR_SLATE)
        badge_bg = COLOR_BLACK if is_lead else COLOR_BORDER
        badge_txt_color = COLOR_BG if is_lead else COLOR_NAVY
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos + Inches(0.12), Inches(6.65), Inches(1.2), Inches(0.22), badge_bg)
        add_text(slide, left_pos + Inches(0.12), Inches(6.67), Inches(1.2), Inches(0.2), badge, FONT_HEADING, 6.8, badge_txt_color, bold=True, align=PP_ALIGN.CENTER)


def build_slide_02_problem(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 1. Background Grid Pattern
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_BG)
    grid_spacing = 0.5
    for x_step in range(1, 26):
        x = x_step * grid_spacing
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(0), Inches(0.008), SLIDE_HEIGHT, COLOR_BORDER)
    for y_step in range(1, 15):
        y = y_step * grid_spacing
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(y), SLIDE_WIDTH, Inches(0.008), COLOR_BORDER)

    # 2. Top Header
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.48), Inches(0.08), Inches(0.30), COLOR_RED)
    add_text(slide, Inches(0.96), Inches(0.50), Inches(6.0), Inches(0.3),
             "THE PROBLEM — 4 SUBSYSTEM FAILURE VOIDS", FONT_MONO, 11, COLOR_BLACK, bold=True)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.48), Inches(2.7), Inches(0.30), RGBColor(0xFF, 0xF1, 0xF2), RGBColor(0xFE, 0xCD, 0xD3))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(9.95), Inches(0.58), Inches(0.08), Inches(0.08), COLOR_RED)
    add_text(slide, Inches(10.1), Inches(0.52), Inches(2.3), Inches(0.22), "FIELD BOTTLENECKS", FONT_MONO, 7.5, COLOR_RED, bold=True, align=PP_ALIGN.CENTER)

    # 3. Compact Headline & Quick Highlights
    add_text(slide, Inches(0.8), Inches(0.95), Inches(7.5), Inches(0.65),
             "Traditional Search & Rescue Fails Across 4 Critical Voids",
             FONT_HEADING, 21, COLOR_BLACK, bold=True)

    quick_badges = ["70% CANOPY CRASHES", "15-30m GEO ERROR", "2-3 HOURS DELAY", "₹40L+ AIRFRAMES"]
    for i, b_text in enumerate(quick_badges):
        bx = Inches(8.3) + i * Inches(1.1)
        is_em = "CRASHES" in b_text or "₹40L" in b_text
        bg_col = RGBColor(0xFF, 0xF1, 0xF2) if is_em else COLOR_CARD_BG
        txt_col = COLOR_RED if is_em else COLOR_SLATE
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(1.05), Inches(1.05), Inches(0.3), bg_col, COLOR_BORDER)
        add_text(slide, bx, Inches(1.10), Inches(1.05), Inches(0.2), b_text, FONT_MONO, 6.0, txt_col, bold=True, align=PP_ALIGN.CENTER)

    # 4. 4-Column Subsystem Problem Void Grid
    problem_voids = [
        ("SUBSYSTEM A // GNC", "Canopy GPS Drift & Crashes", "GPS MULTI-PATH",
         "Dense foliage blocks satellite GNSS signals",
         "Commercial flight controllers lose position lock, drifting into tree canopies and causing catastrophic rotor loss.",
         "70% Loss", "Canopy Sorties", "0 Hold", "Non-GPS Stability", "TARGET: PX4 & 3D ORCA", "LEAD: NIKHIL", COLOR_RED),

        ("SUBSYSTEM B // COMMS", "Mountain Ravine RF Blackout", "DIGITAL CLIFF",
         "Ridgelines sever direct line-of-sight RF",
         "Conventional H.264 digital video completely cuts out below 5dB SNR, plunging base operators into total blindness.",
         "<5dB SNR", "Video Blackout", "0 Relay", "Single-Drone Link", "TARGET: DEEP JSCC MESH", "LEAD: NIKHIL", COLOR_RED),

        ("SUBSYSTEM C // VISION", "Flat-Earth Elevation Drift", "35% FALSE ALARMS",
         "2D raycasts assume flat zero-elevation ground",
         "Sloping terrains produce 15–30m calculated coordinate errors, routing emergency ground teams to empty ravines.",
         "15–30m", "Location Drift", "35%", "False Alarm Rate", "TARGET: 3D DEM RAYCAST", "LEAD: VEDANTH", COLOR_RED),

        ("SUBSYSTEM D // C2 GCS", "Central Pilot Bottleneck", "15-25 CREW LOAD",
         "1-pilot-per-drone manual radio control",
         "Requires 15–25 field personnel and 45–90 min setup; sortie collapses immediately if the single pilot link drops.",
         "2–3 Hrs", "Search Time / mi²", "₹12.5L", "Cost / Deployment", "TARGET: WEBGPU ATAK GCS", "LEAD: SIVA", COLOR_RED),
    ]

    card_w = Inches(2.78)
    card_gap = Inches(0.2)
    for i, (sub_tag, sub_title, sub_badge, void_mech, void_impact, s1_val, s1_lbl, s2_val, s2_lbl, f1, f2, b_col) in enumerate(problem_voids):
        cx = Inches(0.8) + i * (card_w + card_gap)
        cy = Inches(1.68)
        card_h = Inches(4.35)

        # Card Box
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h, COLOR_BG, COLOR_BORDER)

        # Top Badge & Tag
        add_text(slide, cx + Inches(0.12), cy + Inches(0.12), Inches(1.3), Inches(0.2), sub_tag, FONT_MONO, 7.0, COLOR_DIM, bold=True)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(1.45), cy + Inches(0.10), Inches(1.2), Inches(0.22), RGBColor(0xFF, 0xF1, 0xF2), RGBColor(0xFE, 0xCD, 0xD3))
        add_text(slide, cx + Inches(1.45), cy + Inches(0.12), Inches(1.2), Inches(0.2), sub_badge, FONT_MONO, 6.5, COLOR_RED, bold=True, align=PP_ALIGN.CENTER)

        # Title
        add_text(slide, cx + Inches(0.12), cy + Inches(0.38), card_w - Inches(0.24), Inches(0.45), sub_title, FONT_HEADING, 11, COLOR_BLACK, bold=True)

        # Mechanism Box
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.12), cy + Inches(0.90), card_w - Inches(0.24), Inches(0.65), RGBColor(0xFE, 0xF2, 0xF2), RGBColor(0xFE, 0xCD, 0xCD))
        add_text(slide, cx + Inches(0.18), cy + Inches(0.94), card_w - Inches(0.36), Inches(0.18), "CORE FAILURE MECHANISM:", FONT_MONO, 6.0, COLOR_RED, bold=True)
        add_text(slide, cx + Inches(0.18), cy + Inches(1.12), card_w - Inches(0.36), Inches(0.4), void_mech, FONT_BODY, 7.5, COLOR_BLACK, bold=True)

        # Impact Box
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.12), cy + Inches(1.65), card_w - Inches(0.24), Inches(1.45), COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, cx + Inches(0.18), cy + Inches(1.70), card_w - Inches(0.36), Inches(0.18), "OPERATIONAL IMPACT:", FONT_MONO, 6.5, COLOR_BLACK, bold=True)
        add_text(slide, cx + Inches(0.18), cy + Inches(1.92), card_w - Inches(0.36), Inches(1.1), void_impact, FONT_MONO, 7.2, COLOR_SLATE)

        # Specs Grid
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.12), cy + Inches(3.20), Inches(1.22), Inches(0.55), COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, cx + Inches(0.12), cy + Inches(3.24), Inches(1.22), Inches(0.25), s1_val, FONT_MONO, 9.5, COLOR_RED, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, cx + Inches(0.12), cy + Inches(3.52), Inches(1.22), Inches(0.2), s1_lbl, FONT_MONO, 6.0, COLOR_MUTED, align=PP_ALIGN.CENTER)

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(1.44), cy + Inches(3.20), Inches(1.22), Inches(0.55), COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, cx + Inches(1.44), cy + Inches(3.24), Inches(1.22), Inches(0.25), s2_val, FONT_MONO, 9.5, COLOR_BLACK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, cx + Inches(1.44), cy + Inches(3.52), Inches(1.22), Inches(0.2), s2_lbl, FONT_MONO, 6.0, COLOR_MUTED, align=PP_ALIGN.CENTER)

        # Bottom Subsystem Strip
        add_text(slide, cx + Inches(0.12), cy + Inches(3.90), Inches(1.3), Inches(0.2), f1, FONT_MONO, 6.0, COLOR_MUTED)
        add_text(slide, cx + Inches(1.35), cy + Inches(3.90), Inches(1.3), Inches(0.2), f2, FONT_MONO, 6.0, COLOR_RED, bold=True, align=PP_ALIGN.RIGHT)

        # Accent Bar at Bottom
        add_shape(slide, MSO_SHAPE.RECTANGLE, cx, cy + card_h - Inches(0.04), card_w, Inches(0.04), b_col)

    # 5. Bottom Disaster Citations Banner
    disaster_stats = [
        ("FIELD EVIDENCE", "Wayanad Landslide Audit", COLOR_BLACK),
        ("AUDIT FINDING", "70% RF/GPS Disconnect", COLOR_RED),
        ("COST PENALTY", "₹40L Military Airframes", COLOR_BLACK),
        ("SURVIVOR RISK", "412% Disaster Surge", COLOR_RED),
    ]
    banner_w = Inches(2.78)
    banner_gap = Inches(0.2)
    for i, (lbl, val, col) in enumerate(disaster_stats):
        bx = Inches(0.8) + i * (banner_w + banner_gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(6.15), banner_w, Inches(0.48), COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, bx + Inches(0.1), Inches(6.22), Inches(1.3), Inches(0.2), lbl, FONT_MONO, 6.5, COLOR_MUTED)
        add_text(slide, bx + Inches(1.2), Inches(6.20), Inches(1.45), Inches(0.25), val, FONT_MONO, 8.5, col, bold=True, align=PP_ALIGN.RIGHT)

    # 6. Footer
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.75), Inches(11.733), Inches(0.015), COLOR_BORDER)
    add_text(slide, Inches(0.8), Inches(6.85), Inches(3.0), Inches(0.3),
             "PAGE 02", FONT_MONO, 9, COLOR_MUTED, bold=False)
    add_text(slide, Inches(7.5), Inches(6.85), Inches(5.033), Inches(0.3),
             "TEAM OFFGRID — PROJECT SUTRA", FONT_MONO, 9, COLOR_NAVY, bold=True, align=PP_ALIGN.RIGHT)


def build_slide_03_4subsystems(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 1. Background Grid Pattern
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_BG)
    grid_spacing = 0.5
    for x_step in range(1, 26):
        x = x_step * grid_spacing
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(0), Inches(0.008), SLIDE_HEIGHT, COLOR_BORDER)
    for y_step in range(1, 15):
        y = y_step * grid_spacing
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(y), SLIDE_WIDTH, Inches(0.008), COLOR_BORDER)

    # 2. Top Header
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.48), Inches(0.08), Inches(0.30), COLOR_BLACK)
    add_text(slide, Inches(0.96), Inches(0.50), Inches(6.0), Inches(0.3),
             "THE SOLUTION — 4 CORE SUBSYSTEM MOATS", FONT_MONO, 11, COLOR_BLACK, bold=True)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.48), Inches(2.7), Inches(0.30), RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0xA7, 0xF3, 0xD0))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(9.95), Inches(0.58), Inches(0.08), Inches(0.08), COLOR_EMERALD)
    add_text(slide, Inches(10.1), Inches(0.52), Inches(2.3), Inches(0.22), "SWARM ARCHITECTURE", FONT_MONO, 7.5, COLOR_EMERALD, bold=True, align=PP_ALIGN.CENTER)

    # 3. Compact Headline & Quick Highlights
    add_text(slide, Inches(0.8), Inches(0.95), Inches(7.5), Inches(0.65),
             "AI Swarm Drones Find Survivors 3-4x Faster Than Traditional SAR",
             FONT_HEADING, 21, COLOR_BLACK, bold=True)

    quick_badges = ["10-18 MIN/MI²", "78-85% DETECTION", "12% FALSE ALARMS", "₹95,000 SORTIE"]
    for i, b_text in enumerate(quick_badges):
        bx = Inches(8.3) + i * Inches(1.1)
        is_em = "DETECTION" in b_text or "₹95,000" in b_text
        bg_col = RGBColor(0xEC, 0xFD, 0xF5) if is_em else COLOR_CARD_BG
        txt_col = COLOR_EMERALD if is_em else COLOR_SLATE
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(1.05), Inches(1.05), Inches(0.3), bg_col, COLOR_BORDER)
        add_text(slide, bx, Inches(1.10), Inches(1.05), Inches(0.2), b_text, FONT_MONO, 6.0, txt_col, bold=True, align=PP_ALIGN.CENTER)

    # 4. 4-Column Subsystem Architecture Grid
    subsystems = [
        ("SUBSYSTEM A // GNC", "SUTRA-FSD GNC & 3D ORCA", "50Hz OFFBOARD",
         "Canopy GPS Multi-Path Drift & Crashes",
         "Decentralized 3D velocity obstacle collision avoidance + PX4 MicroXRCE 50Hz offboard setpoint streaming.",
         "0.05s", "Latency", "0 Tower", "Reliance", "DECENTRALIZED VIO", "50Hz RATE", COLOR_EMERALD),

        ("SUBSYSTEM B // COMMS", "Deep JSCC Neural Video", "-5dB RESILIENT",
         "Mountain Ravine RF Ridge Blackout",
         "Differentiable joint source-channel coding yielding zero digital cliff and graceful video degradation down to -5dB.",
         "-5dB", "SNR Thresh", "0 Cliff", "Degradation", "SEMANTIC COMMS", "EDGE TPU", COLOR_BLACK),

        ("SUBSYSTEM C // VISION", "3D DEM AI Geolocation", "3.59cm ACCURACY",
         "35% False Alarms & Flat GPS Drift",
         "Jetson Orin YOLOv8-Pose with real-time WGS84 elevation raycasting, eliminating 15-30m calculation errors.",
         "<15ms", "Inference", "3.59cm", "Geo Error", "ELEVATION CORRECTED", "JETSON ORIN", COLOR_EMERALD),

        ("SUBSYSTEM D // C2 GCS", "Pegasus 3D WebGPU GCS", "1-2 OPERATORS",
         "15-25 Ground Crew Cognitive Overload",
         "Offline tactical 3D digital twin + ATAK plugin, empowering 1-2 operators to orchestrate full drone swarms.",
         "1-2 Ops", "Crew Size", "100%", "Offline ATAK", "ATAK INTEGRATION", "WEBGPU 3D", COLOR_BLACK),
    ]

    card_w = Inches(2.78)
    card_gap = Inches(0.2)
    for i, (sub_tag, sub_title, sub_badge, void_text, moat_text, s1_val, s1_lbl, s2_val, s2_lbl, f1, f2, b_col) in enumerate(subsystems):
        cx = Inches(0.8) + i * (card_w + card_gap)
        cy = Inches(1.68)
        card_h = Inches(4.35)

        # Card Box
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h, COLOR_BG, COLOR_BORDER)

        # Top Badge & Tag
        add_text(slide, cx + Inches(0.12), cy + Inches(0.12), Inches(1.3), Inches(0.2), sub_tag, FONT_MONO, 7.0, COLOR_DIM, bold=True)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(1.45), cy + Inches(0.10), Inches(1.2), Inches(0.22), RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0xA7, 0xF3, 0xD0))
        add_text(slide, cx + Inches(1.45), cy + Inches(0.12), Inches(1.2), Inches(0.2), sub_badge, FONT_MONO, 6.5, COLOR_EMERALD, bold=True, align=PP_ALIGN.CENTER)

        # Title
        add_text(slide, cx + Inches(0.12), cy + Inches(0.38), card_w - Inches(0.24), Inches(0.45), sub_title, FONT_HEADING, 11, COLOR_BLACK, bold=True)

        # Problem Void Box
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.12), cy + Inches(0.90), card_w - Inches(0.24), Inches(0.65), RGBColor(0xFE, 0xF2, 0xF2), RGBColor(0xFE, 0xCD, 0xCD))
        add_text(slide, cx + Inches(0.18), cy + Inches(0.94), card_w - Inches(0.36), Inches(0.18), "SOLVES PROBLEM VOID:", FONT_MONO, 6.0, COLOR_RED, bold=True)
        add_text(slide, cx + Inches(0.18), cy + Inches(1.12), card_w - Inches(0.36), Inches(0.4), void_text, FONT_BODY, 7.5, COLOR_BLACK, bold=True)

        # Moat Box
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.12), cy + Inches(1.65), card_w - Inches(0.24), Inches(1.45), COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, cx + Inches(0.18), cy + Inches(1.70), card_w - Inches(0.36), Inches(0.18), "SOLUTIONAL MOAT:", FONT_MONO, 6.5, COLOR_EMERALD, bold=True)
        add_text(slide, cx + Inches(0.18), cy + Inches(1.92), card_w - Inches(0.36), Inches(1.1), moat_text, FONT_MONO, 7.2, COLOR_SLATE)

        # Specs Grid
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.12), cy + Inches(3.20), Inches(1.22), Inches(0.55), COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, cx + Inches(0.12), cy + Inches(3.24), Inches(1.22), Inches(0.25), s1_val, FONT_MONO, 9.5, COLOR_BLACK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, cx + Inches(0.12), cy + Inches(3.52), Inches(1.22), Inches(0.2), s1_lbl, FONT_MONO, 6.0, COLOR_MUTED, align=PP_ALIGN.CENTER)

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(1.44), cy + Inches(3.20), Inches(1.22), Inches(0.55), COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, cx + Inches(1.44), cy + Inches(3.24), Inches(1.22), Inches(0.25), s2_val, FONT_MONO, 9.5, COLOR_EMERALD, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, cx + Inches(1.44), cy + Inches(3.52), Inches(1.22), Inches(0.2), s2_lbl, FONT_MONO, 6.0, COLOR_MUTED, align=PP_ALIGN.CENTER)

        # Bottom Subsystem Strip
        add_text(slide, cx + Inches(0.12), cy + Inches(3.90), Inches(1.3), Inches(0.2), f1, FONT_MONO, 6.0, COLOR_MUTED)
        add_text(slide, cx + Inches(1.35), cy + Inches(3.90), Inches(1.3), Inches(0.2), f2, FONT_MONO, 6.0, COLOR_BLACK, bold=True, align=PP_ALIGN.RIGHT)

        # Accent Bar at Bottom
        add_shape(slide, MSO_SHAPE.RECTANGLE, cx, cy + card_h - Inches(0.04), card_w, Inches(0.04), b_col)

    # 5. Market Validation Banner (Bottom Row)
    market_stats = [
        ("TAM PROJECTION", "₹74,500 Cr Market", COLOR_BLACK),
        ("DEMAND DELTA", "+412% Surge", COLOR_EMERALD),
        ("REGULATORY", "37 Nations Mandate", COLOR_BLACK),
        ("NDMA / SDRF GRANTS", "₹1,500 Cr Allocation", COLOR_EMERALD),
    ]
    banner_w = Inches(2.78)
    banner_gap = Inches(0.2)
    for i, (lbl, val, col) in enumerate(market_stats):
        bx = Inches(0.8) + i * (banner_w + banner_gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(6.15), banner_w, Inches(0.48), COLOR_CARD_BG, COLOR_BORDER)
        add_text(slide, bx + Inches(0.1), Inches(6.22), Inches(1.3), Inches(0.2), lbl, FONT_MONO, 6.5, COLOR_MUTED)
        add_text(slide, bx + Inches(1.2), Inches(6.20), Inches(1.45), Inches(0.25), val, FONT_MONO, 8.5, col, bold=True, align=PP_ALIGN.RIGHT)

    # 6. Footer
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.75), Inches(11.733), Inches(0.015), COLOR_BORDER)
    add_text(slide, Inches(0.8), Inches(6.85), Inches(3.0), Inches(0.3),
             "PAGE 03", FONT_MONO, 9, COLOR_MUTED, bold=False)
    add_text(slide, Inches(7.5), Inches(6.85), Inches(5.033), Inches(0.3),
             "TEAM OFFGRID — PROJECT SUTRA", FONT_MONO, 9, COLOR_NAVY, bold=True, align=PP_ALIGN.RIGHT)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    prs.core_properties.title = "PROJECT SUTRA Pitch Deck"
    prs.core_properties.author = "Team Offgrid"

    print("Building Slide 1 (Title), Slide 2 (4 Subsystem Problem Voids), & Slide 3 (4 Subsystem Moats)...")
    build_slide_01_title(prs)
    build_slide_02_problem(prs)
    build_slide_03_4subsystems(prs)

    output_path = Path("sutra_pitch_deck.pptx")
    prs.save(str(output_path))
    print(f"✓ Presentation saved to: {output_path.resolve()}")


if __name__ == "__main__":
    main()
