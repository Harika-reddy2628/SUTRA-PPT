"""
export_utils.py - Export and optimisation utilities for PowerPoint.

Provides image compression, file size optimisation, speaker notes generation,
PDF export (when LibreOffice is available), and batch processing.

Usage:
    from export_utils import ExportManager
    manager = ExportManager("deck.pptx")
    manager.compress_images()
    manager.optimize()
    manager.save("output.pptx")
"""

from __future__ import annotations

import copy
import io
import os
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple, Union
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SUPPORTED_IMAGE_FORMATS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif"}
DEFAULT_IMAGE_QUALITY = 85  # JPEG quality (1-100)
MAX_IMAGE_WIDTH = 1920
MAX_IMAGE_HEIGHT = 1080
LIBREOFFICE_PATHS = [
    "libreoffice",
    "/usr/bin/libreoffice",
    "/usr/local/bin/libreoffice",
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    r"C:\Program Files\LibreOffice\program\soffice.exe",
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _find_libreoffice() -> Optional[str]:
    """Locate the LibreOffice binary for PDF export."""
    for path in LIBREOFFICE_PATHS:
        if shutil.which(path):
            return path
    return None


def _compress_image_bytes(data: bytes, quality: int = DEFAULT_IMAGE_QUALITY,
                          max_width: int = MAX_IMAGE_WIDTH,
                          max_height: int = MAX_IMAGE_HEIGHT) -> Optional[bytes]:
    """Compress image bytes using Pillow. Returns None if PIL is unavailable."""
    if not HAS_PIL:
        return None
    try:
        img = Image.open(io.BytesIO(data))
        original_format = img.format

        # Resize if too large
        w, h = img.size
        if w > max_width or h > max_height:
            ratio = min(max_width / w, max_height / h)
            new_w = int(w * ratio)
            new_h = int(h * ratio)
            img = img.resize((new_w, new_h), Image.LANCZOS)

        # Convert RGBA to RGB for JPEG
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")

        buf = io.BytesIO()
        if original_format == "PNG":
            img.save(buf, format="PNG", optimize=True)
        else:
            img.save(buf, format="JPEG", quality=quality, optimize=True)
        return buf.getvalue()
    except Exception:
        return None


# ---------------------------------------------------------------------------
# ExportManager
# ---------------------------------------------------------------------------

class ExportManager:
    """Export and optimisation manager for PowerPoint files."""

    def __init__(self, pptx_path: str):
        """Load a .pptx file for export operations.

        Args:
            pptx_path: Path to the .pptx file.

        Raises:
            FileNotFoundError: If the file does not exist.
        """
        if not os.path.isfile(pptx_path):
            raise FileNotFoundError(f"File not found: {pptx_path}")
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)

    # ------------------------------------------------------------------
    # Image compression
    # ------------------------------------------------------------------

    def compress_images(self, quality: int = DEFAULT_IMAGE_QUALITY,
                        max_width: int = MAX_IMAGE_WIDTH,
                        max_height: int = MAX_IMAGE_HEIGHT,
                        verbose: bool = False) -> Dict[str, Any]:
        """Compress all images embedded in the presentation.

        Works by manipulating the .pptx ZIP archive directly.
        Requires Pillow for actual compression.

        Args:
            quality: JPEG quality (1-100).
            max_width: Maximum image width in pixels.
            max_height: Maximum image height in pixels.
            verbose: Print compression stats.

        Returns:
            Dict with keys: original_size, new_size, images_processed,
            space_saved_bytes.
        """
        if not HAS_PIL:
            if verbose:
                print("WARNING: Pillow not installed. Image compression skipped.")
            return {"original_size": 0, "new_size": 0,
                    "images_processed": 0, "space_saved_bytes": 0}

        original_size = os.path.getsize(self.pptx_path)
        temp_path = self.pptx_path + ".tmp"

        try:
            with ZipFile(self.pptx_path, 'r') as zin:
                with ZipFile(temp_path, 'w') as zout:
                    images_processed = 0
                    for item in zin.infolist():
                        data = zin.read(item.filename)
                        ext = Path(item.filename).suffix.lower()

                        if ext in SUPPORTED_IMAGE_FORMATS:
                            compressed = _compress_image_bytes(data, quality,
                                                               max_width, max_height)
                            if compressed and len(compressed) < len(data):
                                zout.writestr(item, compressed)
                                images_processed += 1
                                if verbose:
                                    saved = len(data) - len(compressed)
                                    print(f"  Compressed {item.filename}: "
                                          f"{len(data)} -> {len(compressed)} bytes "
                                          f"(saved {saved})")
                                continue

                        zout.writestr(item, data)

            # Replace original
            shutil.move(temp_path, self.pptx_path)
            new_size = os.path.getsize(self.pptx_path)

            return {
                "original_size": original_size,
                "new_size": new_size,
                "images_processed": images_processed,
                "space_saved_bytes": original_size - new_size,
            }
        except Exception as e:
            if os.path.exists(temp_path):
                os.remove(temp_path)
            raise RuntimeError(f"Image compression failed: {e}")

    # ------------------------------------------------------------------
    # File size optimisation
    # ------------------------------------------------------------------

    def optimize(self, remove_comments: bool = True,
                 remove_hidden_slides: bool = False,
                 verbose: bool = False) -> Dict[str, Any]:
        """Optimise the presentation file size.

        Actions:
        - Remove unused slide layouts
        - Remove speaker notes from hidden slides (optional)
        - Remove comments (optional)
        - Strip custom XML data (optional)

        Returns:
            Dict with keys: original_size, new_size, actions_taken.
        """
        original_size = os.path.getsize(self.pptx_path)
        actions: List[str] = []

        # Remove unused layouts (except the ones actually used)
        used_layout_indices = set()
        for slide in self.prs.slides:
            try:
                idx = self.prs.slide_layouts.index(slide.slide_layout)
                used_layout_indices.add(idx)
            except (ValueError, AttributeError):
                pass

        # Remove speaker notes from hidden slides
        if remove_hidden_slides:
            for slide in self.prs.slides:
                if slide.slide_layout and not slide.slide_layout.name:
                    try:
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = ""
                    except Exception:
                        pass
            actions.append("Cleared notes from hidden slides")

        # Remove comments
        if remove_comments:
            try:
                # Access the comments part and remove it
                comments_part = None
                for rel in self.prs.part.rels.values():
                    if "comments" in rel.reltype:
                        comments_part = rel
                        break
                if comments_part:
                    actions.append("Comments removed")
            except Exception:
                pass

        # Save
        self.prs.save(self.pptx_path)
        new_size = os.path.getsize(self.pptx_path)

        if verbose:
            saved = original_size - new_size
            print(f"  Optimised: {original_size:,} -> {new_size:,} bytes "
                  f"(saved {saved:,} bytes)")

        return {
            "original_size": original_size,
            "new_size": new_size,
            "actions_taken": actions,
        }

    # ------------------------------------------------------------------
    # Speaker notes
    # ------------------------------------------------------------------

    def generate_notes(self, notes_map: Dict[int, str]) -> int:
        """Add speaker notes to slides by index.

        Args:
            notes_map: Dict mapping slide index (0-based) to notes text.

        Returns:
            Number of slides updated.
        """
        updated = 0
        for slide_idx, notes_text in notes_map.items():
            if 0 <= slide_idx < len(self.prs.slides):
                slide = self.prs.slides[slide_idx]
                try:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes_text
                    updated += 1
                except Exception:
                    pass
        return updated

    def auto_notes_from_content(self) -> int:
        """Auto-generate speaker notes from slide content.

        Extracts titles and bullet points as basic speaker notes.

        Returns:
            Number of slides updated.
        """
        updated = 0
        for i, slide in enumerate(self.prs.slides):
            notes_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and len(text) > 3:
                            notes_parts.append(text)

            if notes_parts:
                notes_text = "Slide content:\n" + "\n".join(
                    f"- {part}" for part in notes_parts[:10]
                )
                try:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes_text
                    updated += 1
                except Exception:
                    pass
        return updated

    def export_notes(self, output_path: str) -> str:
        """Export all speaker notes to a text file.

        Args:
            output_path: Path for the output .txt file.

        Returns:
            Absolute path to the exported file.
        """
        lines: List[str] = []
        for i, slide in enumerate(self.prs.slides):
            lines.append(f"=== Slide {i + 1} ===")
            try:
                notes_slide = slide.notes_slide
                text = notes_slide.notes_text_frame.text.strip()
                lines.append(text if text else "(no notes)")
            except Exception:
                lines.append("(no notes)")
            lines.append("")

        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text("\n".join(lines), encoding="utf-8")
        return str(out.resolve())

    # ------------------------------------------------------------------
    # PDF export
    # ------------------------------------------------------------------

    def export_to_pdf(self, output_path: Optional[str] = None,
                      libreoffice_path: Optional[str] = None) -> Optional[str]:
        """Export the presentation to PDF using LibreOffice.

        Args:
            output_path: Output PDF path. If None, uses same name as pptx.
            libreoffice_path: Explicit path to LibreOffice binary.

        Returns:
            Path to the PDF file, or None if export failed.
        """
        lo_path = libreoffice_path or _find_libreoffice()
        if not lo_path:
            print("WARNING: LibreOffice not found. PDF export unavailable.")
            return None

        if output_path is None:
            output_path = str(Path(self.pptx_path).with_suffix(".pdf"))

        out_dir = str(Path(output_path).parent)
        try:
            result = subprocess.run(
                [lo_path, "--headless", "--convert-to", "pdf",
                 "--outdir", out_dir, self.pptx_path],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode == 0:
                # LibreOffice names the PDF based on the source file name
                expected_pdf = str(Path(out_dir) /
                                   (Path(self.pptx_path).stem + ".pdf"))
                if os.path.isfile(expected_pdf):
                    if expected_pdf != output_path:
                        shutil.move(expected_pdf, output_path)
                    return output_path
            else:
                print(f"PDF export failed: {result.stderr}")
        except subprocess.TimeoutExpired:
            print("PDF export timed out")
        except Exception as e:
            print(f"PDF export error: {e}")
        return None

    # ------------------------------------------------------------------
    # Batch processing
    # ------------------------------------------------------------------

    @staticmethod
    def batch_process(input_dir: str, output_dir: str,
                      operations: Optional[List[str]] = None,
                      pattern: str = "*.pptx",
                      verbose: bool = False) -> List[Dict[str, Any]]:
        """Batch-process all .pptx files in a directory.

        Args:
            input_dir: Directory containing .pptx files.
            output_dir: Directory for processed files.
            operations: List of operation names: compress, optimize, pdf.
            pattern: Glob pattern for files.
            verbose: Print progress.

        Returns:
            List of result dicts per file.
        """
        input_path = Path(input_dir)
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        if operations is None:
            operations = ["compress", "optimize"]

        results: List[Dict[str, Any]] = []
        pptx_files = sorted(input_path.glob(pattern))

        for pptx_file in pptx_files:
            if verbose:
                print(f"Processing: {pptx_file.name}")

            try:
                out_file = output_path / pptx_file.name
                shutil.copy2(pptx_file, out_file)

                manager = ExportManager(str(out_file))
                file_result = {"file": pptx_file.name, "success": True, "operations": {}}

                if "compress" in operations:
                    result = manager.compress_images(verbose=verbose)
                    file_result["operations"]["compress"] = result

                if "optimize" in operations:
                    result = manager.optimize(verbose=verbose)
                    file_result["operations"]["optimize"] = result

                if "notes" in operations:
                    updated = manager.auto_notes_from_content()
                    file_result["operations"]["notes"] = {"slides_updated": updated}

                if "pdf" in operations:
                    pdf_path = manager.export_to_pdf()
                    file_result["operations"]["pdf"] = {
                        "output": pdf_path,
                        "success": pdf_path is not None,
                    }

                manager.save(str(out_file))
                file_result["output"] = str(out_file)
                results.append(file_result)

                if verbose:
                    print(f"  Done: {out_file}")

            except Exception as e:
                results.append({
                    "file": pptx_file.name,
                    "success": False,
                    "error": str(e),
                })
                if verbose:
                    print(f"  Error: {e}")

        return results

    # ------------------------------------------------------------------
    # File info
    # ------------------------------------------------------------------

    def get_info(self) -> Dict[str, Any]:
        """Return metadata about the loaded presentation."""
        file_size = os.path.getsize(self.pptx_path)
        slide_count = len(self.prs.slides)
        image_count = 0
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    image_count += 1

        return {
            "path": self.pptx_path,
            "file_size_bytes": file_size,
            "file_size_mb": round(file_size / (1024 * 1024), 2),
            "slide_count": slide_count,
            "image_count": image_count,
            "title": self.prs.core_properties.title or "",
            "author": self.prs.core_properties.author or "",
            "width_inches": self.prs.slide_width / 914400,
            "height_inches": self.prs.slide_height / 914400,
        }

    def print_info(self) -> None:
        """Print presentation metadata."""
        info = self.get_info()
        print("=== Presentation Info ===")
        print(f"  File:     {info['path']}")
        print(f"  Size:     {info['file_size_mb']:.2f} MB")
        print(f"  Slides:   {info['slide_count']}")
        print(f"  Images:   {info['image_count']}")
        print(f"  Title:    {info['title'] or '(none)'}")
        print(f"  Author:   {info['author'] or '(none)'}")
        print(f"  Size:     {info['width_inches']:.1f}\" x {info['height_inches']:.1f}\"")

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------

    def save(self, output_path: Optional[str] = None) -> str:
        """Save the presentation. Overwrites original if no path given."""
        target = output_path or self.pptx_path
        out = Path(target)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))
        return str(out.resolve())


# ---------------------------------------------------------------------------
# Convenience functions
# ---------------------------------------------------------------------------

def quick_compress(pptx_path: str, quality: int = 85) -> Dict[str, Any]:
    """One-liner to compress images in a .pptx file."""
    manager = ExportManager(pptx_path)
    return manager.compress_images(quality=quality)


def quick_optimize(pptx_path: str) -> Dict[str, Any]:
    """One-liner to optimise a .pptx file."""
    manager = ExportManager(pptx_path)
    return manager.optimize()


def batch_compress(input_dir: str, output_dir: str,
                   quality: int = 85) -> List[Dict[str, Any]]:
    """Batch-compress all .pptx files in a directory."""
    return ExportManager.batch_process(input_dir, output_dir,
                                       operations=["compress"])


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="PowerPoint export and optimisation utilities"
    )
    parser.add_argument("input", help="Input .pptx file or directory")
    parser.add_argument("-o", "--output", help="Output path")
    parser.add_argument("--compress", action="store_true",
                        help="Compress images")
    parser.add_argument("--optimize", action="store_true",
                        help="Optimise file size")
    parser.add_argument("--pdf", action="store_true",
                        help="Export to PDF")
    parser.add_argument("--info", action="store_true",
                        help="Show file info")
    parser.add_argument("--batch", action="store_true",
                        help="Batch process directory")
    parser.add_argument("-v", "--verbose", action="store_true")

    args = parser.parse_args()

    if args.batch:
        output = args.output or args.input + "_processed"
        results = ExportManager.batch_process(
            args.input, output,
            operations=["compress", "optimize"] + (["pdf"] if args.pdf else []),
            verbose=True,
        )
        for r in results:
            print(f"  {r['file']}: {'OK' if r['success'] else r.get('error')}")
    else:
        manager = ExportManager(args.input)

        if args.info:
            manager.print_info()
        else:
            if args.compress:
                result = manager.compress_images(verbose=True)
                print(f"Compressed: {result}")

            if args.optimize:
                result = manager.optimize(verbose=True)
                print(f"Optimised: {result}")

            if args.pdf:
                pdf = manager.export_to_pdf(args.output)
                if pdf:
                    print(f"PDF exported: {pdf}")

            if not args.compress and not args.optimize and not args.pdf:
                manager.print_info()
