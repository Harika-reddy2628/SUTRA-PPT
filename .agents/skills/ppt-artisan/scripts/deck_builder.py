"""
deck_builder.py - Core PowerPoint deck creation engine.

Provides SlideBuilder class with methods for every slide type,
design system constants, template loading, and theme application.

Usage:
    from deck_builder import SlideBuilder
    builder = SlideBuilder()
    builder.create_deck("My Presentation", "Author Name")
    builder.add_title_slide("Welcome", "Subtitle here")
    builder.save("output.pptx")
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Design System Constants
# ---------------------------------------------------------------------------

COLOR_PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)    # Deep navy
COLOR_SECONDARY = RGBColor(0x2E, 0x86, 0xAB)  # Teal
COLOR_ACCENT = RGBColor(0xE8, 0x4D, 0x3D)     # Coral
COLOR_SUCCESS = RGBColor(0x27, 0xAE, 0x60)    # Green
COLOR_WARNING = RGBColor(0xF3, 0x9C, 0x12)    # Amber
COLOR_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)      # Light gray
COLOR_DARK = RGBColor(0x2C, 0x3E, 0x50)       # Dark gray
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)

FONT_HEADING = "Montserrat"
FONT_BODY = "Inter"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.5)

CHART_COLORS = [COLOR_PRIMARY, COLOR_SECONDARY, COLOR_ACCENT, COLOR_SUCCESS, COLOR_WARNING]


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a hex colour string to RGBColor."""
    hex_str = hex_str.lstrip("#")
    if len(hex_str) != 6:
        raise ValueError(f"Invalid hex colour: {hex_str}")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _set_run_font(run, font_name: str = FONT_BODY, size: int = 18,
                  color: RGBColor = COLOR_DARK, bold: bool = False,
                  italic: bool = False) -> None:
    """Apply font settings to a run (run-level formatting)."""
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def _add_textbox(slide, left: Inches, top: Inches, width: Inches,
                 height: Inches, text: str, font_name: str = FONT_BODY,
                 font_size: int = 18, color: RGBColor = COLOR_DARK,
                 bold: bool = False, alignment: PP_ALIGN = PP_ALIGN.LEFT,
                 anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> Any:
    """Add a styled textbox and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        txBox.text_frame.paragraphs[0].alignment = alignment
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    for run in p.runs:
        _set_run_font(run, font_name, font_size, color, bold)
    return tf


def _add_bullet_list(tf, items: List[str], font_name: str = FONT_BODY,
                     font_size: int = 16, color: RGBColor = COLOR_DARK,
                     level: int = 0, spacing_after: int = 6) -> None:
    """Populate a text frame with bullet items."""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = level
        p.space_after = Pt(spacing_after)
        for run in p.runs:
            _set_run_font(run, font_name, font_size, color)


def _add_shape_rect(slide, left, top, width, height,
                    fill_color: RGBColor, line_color: Optional[RGBColor] = None) -> Any:
    """Add a rectangle shape with optional fill."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# SlideBuilder
# ---------------------------------------------------------------------------

class SlideBuilder:
    """Core PowerPoint deck builder with methods for every slide type."""

    def __init__(self, template_path: Optional[str] = None):
        """Initialise the builder, optionally loading a .pptx template."""
        if template_path and os.path.isfile(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.theme: Dict[str, Any] = {
            "primary": COLOR_PRIMARY,
            "secondary": COLOR_SECONDARY,
            "accent": COLOR_ACCENT,
            "success": COLOR_SUCCESS,
            "warning": COLOR_WARNING,
            "heading_font": FONT_HEADING,
            "body_font": FONT_BODY,
        }

    # ------------------------------------------------------------------
    # Metadata
    # ------------------------------------------------------------------

    def create_deck(self, title: str, author: str,
                    theme: Optional[Dict[str, Any]] = None) -> "SlideBuilder":
        """Set deck metadata and optionally override the theme."""
        if theme:
            self.theme.update(theme)
        self.prs.core_properties.title = title
        self.prs.core_properties.author = author
        return self

    def apply_theme(self, theme: Dict[str, Any]) -> "SlideBuilder":
        """Apply a theme dictionary to override default colours/fonts."""
        self.theme.update(theme)
        return self

    # ------------------------------------------------------------------
    # Slide helpers
    # ------------------------------------------------------------------

    def _blank_slide(self):
        """Add a blank slide layout."""
        layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(layout)

    # ------------------------------------------------------------------
    # Slide methods
    # ------------------------------------------------------------------

    def add_title_slide(self, title: str, subtitle: str = "",
                        bg_color: Optional[RGBColor] = None) -> "SlideBuilder":
        """Add a title slide with large heading and subtitle."""
        bg = bg_color or self.theme["primary"]
        slide = self._blank_slide()

        # Background
        _add_shape_rect(slide, Inches(0), Inches(0),
                        SLIDE_WIDTH, SLIDE_HEIGHT, bg)

        # Title
        _add_textbox(slide, Inches(1), Inches(2.2), Inches(11.333), Inches(2),
                     title, self.theme["heading_font"], 44, COLOR_WHITE, True,
                     PP_ALIGN.CENTER)

        # Subtitle
        if subtitle:
            _add_textbox(slide, Inches(1), Inches(4.4), Inches(11.333), Inches(1.2),
                         subtitle, self.theme["body_font"], 24, COLOR_LIGHT,
                         False, PP_ALIGN.CENTER)

        # Accent line
        line_w = Inches(3)
        _add_shape_rect(slide, Inches(5.166), Inches(4.1), line_w, Inches(0.06),
                        self.theme["accent"])
        return self

    def add_content_slide(self, title: str, bullets: List[str],
                          accent_color: Optional[RGBColor] = None) -> "SlideBuilder":
        """Add a content slide with title and bullet list."""
        accent = accent_color or self.theme["secondary"]
        slide = self._blank_slide()

        # Top accent bar
        _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), accent)

        # Title
        _add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.9),
                     title, self.theme["heading_font"], 32, COLOR_PRIMARY, True,
                     PP_ALIGN.LEFT)

        # Divider line
        _add_shape_rect(slide, Inches(0.8), Inches(1.45), Inches(2), Inches(0.05),
                        accent)

        # Bullet list
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        _add_bullet_list(tf, bullets, self.theme["body_font"], 20, COLOR_DARK)
        return self

    def add_section_divider(self, title: str, subtitle: str = "") -> "SlideBuilder":
        """Add a bold section divider slide."""
        slide = self._blank_slide()
        _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
                        self.theme["primary"])

        # Section number / decorative element
        _add_shape_rect(slide, Inches(1), Inches(2.5), Inches(0.8), Inches(0.06),
                        self.theme["accent"])

        _add_textbox(slide, Inches(1), Inches(2.8), Inches(11.333), Inches(1.5),
                     title, self.theme["heading_font"], 40, COLOR_WHITE, True,
                     PP_ALIGN.LEFT)

        if subtitle:
            _add_textbox(slide, Inches(1), Inches(4.4), Inches(11.333), Inches(0.8),
                         subtitle, self.theme["body_font"], 20, COLOR_LIGHT,
                         False, PP_ALIGN.LEFT)
        return self

    def add_comparison(self, title: str,
                       left_title: str, left_items: List[str],
                       right_title: str, right_items: List[str]) -> "SlideBuilder":
        """Add a side-by-side comparison slide."""
        slide = self._blank_slide()

        # Title
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9),
                     title, self.theme["heading_font"], 32, COLOR_PRIMARY, True,
                     PP_ALIGN.LEFT)
        _add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05),
                        self.theme["accent"])

        half_w = Inches(5.4)

        # Left column
        _add_shape_rect(slide, Inches(0.8), Inches(1.7), half_w, Inches(5.2),
                        COLOR_LIGHT)
        _add_textbox(slide, Inches(1.2), Inches(1.9), Inches(4.6), Inches(0.6),
                     left_title, self.theme["heading_font"], 22,
                     self.theme["primary"], True, PP_ALIGN.LEFT)

        txBox_l = slide.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(4.6), Inches(4))
        tf_l = txBox_l.text_frame
        tf_l.word_wrap = True
        _add_bullet_list(tf_l, left_items, self.theme["body_font"], 16, COLOR_DARK)

        # Right column
        _add_shape_rect(slide, Inches(6.9), Inches(1.7), half_w, Inches(5.2),
                        COLOR_LIGHT)
        _add_textbox(slide, Inches(7.3), Inches(1.9), Inches(4.6), Inches(0.6),
                     right_title, self.theme["heading_font"], 22,
                     self.theme["accent"], True, PP_ALIGN.LEFT)

        txBox_r = slide.shapes.add_textbox(Inches(7.3), Inches(2.6), Inches(4.6), Inches(4))
        tf_r = txBox_r.text_frame
        tf_r.word_wrap = True
        _add_bullet_list(tf_r, right_items, self.theme["body_font"], 16, COLOR_DARK)

        # Centre divider
        _add_shape_rect(slide, Inches(6.45), Inches(1.7), Inches(0.04), Inches(5.2),
                        self.theme["accent"])
        return self

    def add_timeline(self, title: str,
                     events: List[Tuple[str, str, str]]) -> "SlideBuilder":
        """Add a timeline slide. Events: [(date, label, description), ...]"""
        slide = self._blank_slide()

        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9),
                     title, self.theme["heading_font"], 32, COLOR_PRIMARY, True,
                     PP_ALIGN.LEFT)
        _add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05),
                        self.theme["accent"])

        n = len(events)
        if n == 0:
            return self

        # Horizontal line
        line_y = Inches(3.2)
        _add_shape_rect(slide, Inches(0.8), line_y, Inches(11.733), Inches(0.04),
                        self.theme["secondary"])

        step = 11.733 / max(n, 1)
        for i, (date, label, desc) in enumerate(events):
            cx = Inches(0.8 + step * i + step / 2)

            # Dot
            from pptx.enum.shapes import MSO_SHAPE
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         cx - Inches(0.12), line_y - Inches(0.1),
                                         Inches(0.24), Inches(0.24))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.theme["accent"]
            dot.line.fill.background()

            # Date (above)
            _add_textbox(slide, cx - Inches(0.9), line_y - Inches(0.7),
                         Inches(1.8), Inches(0.5), date,
                         self.theme["body_font"], 12, self.theme["secondary"],
                         True, PP_ALIGN.CENTER)

            # Label (below)
            _add_textbox(slide, cx - Inches(0.9), line_y + Inches(0.3),
                         Inches(1.8), Inches(0.5), label,
                         self.theme["heading_font"], 14, COLOR_PRIMARY,
                         True, PP_ALIGN.CENTER)

            # Description
            _add_textbox(slide, cx - Inches(0.9), line_y + Inches(0.8),
                         Inches(1.8), Inches(1.2), desc,
                         self.theme["body_font"], 11, COLOR_DARK,
                         False, PP_ALIGN.CENTER)
        return self

    def add_chart_slide(self, title: str, chart_type: str,
                        data: List[List[float]], labels: List[str],
                        series_names: Optional[List[str]] = None,
                        colors: Optional[List[RGBColor]] = None) -> "SlideBuilder":
        """Add a chart slide. chart_type: bar, line, pie."""
        slide = self._blank_slide()
        colors = colors or CHART_COLORS

        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9),
                     title, self.theme["heading_font"], 32, COLOR_PRIMARY, True,
                     PP_ALIGN.LEFT)
        _add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05),
                        self.theme["accent"])

        chart_data = ChartData()
        chart_data.categories = labels

        if series_names is None:
            series_names = [f"Series {i+1}" for i in range(len(data))]
        for name, values in zip(series_names, data):
            chart_data.add_series(name, values)

        chart_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        xl_type = chart_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_frame = slide.shapes.add_chart(
            xl_type, Inches(1), Inches(1.7), Inches(11.333), Inches(5.3), chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = len(data) > 1

        # Apply accent colour to first series
        if chart.series:
            try:
                chart.series[0].format.fill.solid()
                chart.series[0].format.fill.fore_color.rgb = self.theme["accent"]
            except Exception:
                pass
        return self

    def add_team_slide(self, title: str,
                       members: List[Tuple[str, str, Optional[str]]]) -> "SlideBuilder":
        """Add a team slide. Members: [(name, role, photo_path), ...]"""
        slide = self._blank_slide()

        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9),
                     title, self.theme["heading_font"], 32, COLOR_PRIMARY, True,
                     PP_ALIGN.LEFT)
        _add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05),
                        self.theme["accent"])

        n = len(members)
        if n == 0:
            return self

        cols = min(n, 4)
        card_w = Inches(2.8)
        gap = (11.733 - cols * 2.8) / (cols + 1) if cols > 0 else 0

        for i, (name, role, photo_path) in enumerate(members):
            col = i % cols
            row = i // cols
            cx = Inches(0.8 + gap * (col + 1) + 2.8 * col)
            cy = Inches(1.8 + row * 3.2)

            # Card background
            _add_shape_rect(slide, cx, cy, card_w, Inches(2.8), COLOR_LIGHT)

            # Photo placeholder or actual photo
            if photo_path and os.path.isfile(photo_path):
                slide.shapes.add_picture(photo_path, cx + Inches(0.6), cy + Inches(0.2),
                                         Inches(1.6), Inches(1.6))
            else:
                from pptx.enum.shapes import MSO_SHAPE
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                                cx + Inches(0.6), cy + Inches(0.2),
                                                Inches(1.6), Inches(1.6))
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.theme["secondary"]
                circle.line.fill.background()
                _add_textbox(slide, cx + Inches(0.6), cy + Inches(0.65),
                             Inches(1.6), Inches(0.7), name[:2].upper(),
                             self.theme["heading_font"], 28, COLOR_WHITE, True,
                             PP_ALIGN.CENTER)

            # Name
            _add_textbox(slide, cx + Inches(0.1), cy + Inches(1.9),
                         Inches(2.6), Inches(0.5), name,
                         self.theme["heading_font"], 16, COLOR_PRIMARY, True,
                         PP_ALIGN.CENTER)
            # Role
            _add_textbox(slide, cx + Inches(0.1), cy + Inches(2.35),
                         Inches(2.6), Inches(0.4), role,
                         self.theme["body_font"], 13, self.theme["secondary"],
                         False, PP_ALIGN.CENTER)
        return self

    def add_quote_slide(self, quote: str, author: str,
                        role: str = "") -> "SlideBuilder":
        """Add a quote / testimonial slide."""
        slide = self._blank_slide()
        _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
                        self.theme["primary"])

        # Large quotation mark
        _add_textbox(slide, Inches(1), Inches(1.2), Inches(2), Inches(1.5),
                     "\u201C", self.theme["heading_font"], 120, self.theme["accent"],
                     True, PP_ALIGN.LEFT)

        # Quote text
        _add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10.333), Inches(2.5),
                     quote, self.theme["heading_font"], 28, COLOR_WHITE, False,
                     PP_ALIGN.CENTER)

        # Author
        _add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10.333), Inches(0.6),
                     f"\u2014 {author}", self.theme["body_font"], 18,
                     self.theme["accent"], True, PP_ALIGN.CENTER)

        if role:
            _add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.333), Inches(0.5),
                         role, self.theme["body_font"], 14, COLOR_LIGHT, False,
                         PP_ALIGN.CENTER)
        return self

    def add_kpi_dashboard(self, title: str,
                          kpis: List[Tuple[str, str, Optional[str]]]) -> "SlideBuilder":
        """Add a KPI dashboard slide. KPIs: [(label, value, trend), ...]"""
        slide = self._blank_slide()

        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9),
                     title, self.theme["heading_font"], 32, COLOR_PRIMARY, True,
                     PP_ALIGN.LEFT)
        _add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05),
                        self.theme["accent"])

        n = len(kpis)
        if n == 0:
            return self

        cols = min(n, 4)
        card_w = Inches(2.6)
        gap = (11.733 - cols * 2.6) / (cols + 1) if cols > 0 else 0

        trend_colors = {
            "up": COLOR_SUCCESS, "down": COLOR_ACCENT,
            "+": COLOR_SUCCESS, "-": COLOR_ACCENT,
        }

        for i, (label, value, trend) in enumerate(kpis):
            col = i % cols
            row = i // cols
            cx = Inches(0.8 + gap * (col + 1) + 2.6 * col)
            cy = Inches(1.8 + row * 2.8)

            _add_shape_rect(slide, cx, cy, card_w, Inches(2.4), COLOR_LIGHT)

            # Value
            _add_textbox(slide, cx + Inches(0.2), cy + Inches(0.3),
                         Inches(2.2), Inches(1), value,
                         self.theme["heading_font"], 36, self.theme["primary"],
                         True, PP_ALIGN.CENTER)

            # Label
            _add_textbox(slide, cx + Inches(0.2), cy + Inches(1.3),
                         Inches(2.2), Inches(0.5), label,
                         self.theme["body_font"], 14, COLOR_DARK, False,
                         PP_ALIGN.CENTER)

            # Trend indicator
            if trend:
                tc = COLOR_DARK
                for key, color in trend_colors.items():
                    if trend.startswith(key):
                        tc = color
                        break
                _add_textbox(slide, cx + Inches(0.2), cy + Inches(1.8),
                             Inches(2.2), Inches(0.4), trend,
                             self.theme["body_font"], 13, tc, True,
                             PP_ALIGN.CENTER)
        return self

    def add_image_slide(self, title: str, image_path: str,
                        caption: str = "") -> "SlideBuilder":
        """Add a slide with a full-width image and optional caption."""
        slide = self._blank_slide()

        _add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.733), Inches(0.7),
                     title, self.theme["heading_font"], 28, COLOR_PRIMARY, True,
                     PP_ALIGN.LEFT)
        _add_shape_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.05),
                        self.theme["accent"])

        img_top = Inches(1.3)
        img_h = Inches(5.3) if caption else Inches(5.8)

        if os.path.isfile(image_path):
            slide.shapes.add_picture(image_path, Inches(1), img_top,
                                     Inches(11.333), img_h)
        else:
            _add_shape_rect(slide, Inches(1), img_top, Inches(11.333), img_h,
                            COLOR_LIGHT)
            _add_textbox(slide, Inches(1), img_top + Inches(2),
                         Inches(11.333), Inches(1), "[Image not found]",
                         self.theme["body_font"], 18, COLOR_DARK, False,
                         PP_ALIGN.CENTER)

        if caption:
            _add_textbox(slide, Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.5),
                         caption, self.theme["body_font"], 12, self.theme["secondary"],
                         False, PP_ALIGN.CENTER)
        return self

    # ------------------------------------------------------------------
    # Speaker notes
    # ------------------------------------------------------------------

    def add_notes(self, slide_index: int, notes_text: str) -> "SlideBuilder":
        """Add speaker notes to a slide by index (0-based)."""
        if 0 <= slide_index < len(self.prs.slides):
            slide = self.prs.slides[slide_index]
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
        return self

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------

    def save(self, output_path: str) -> str:
        """Save the presentation to disk and return the absolute path."""
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))
        return str(out.resolve())

    # ------------------------------------------------------------------
    # Convenience
    # ------------------------------------------------------------------

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)

    def get_presentation(self) -> Presentation:
        return self.prs


# ---------------------------------------------------------------------------
# CLI entry point for quick testing
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    builder = SlideBuilder()
    builder.create_deck("Demo Deck", "ppt-artisan")
    builder.add_title_slide("Welcome", "Automated presentation generation")
    builder.add_content_slide("Overview", [
        "Bullet one",
        "Bullet two",
        "Bullet three",
    ])
    builder.add_section_divider("Section 1", "Going deeper")
    builder.add_comparison("Pros vs Cons", "Pros", ["Fast", "Cheap"], "Cons", ["Complex"])
    builder.add_timeline("Milestones", [
        ("Q1", "Kickoff", "Project start"),
        ("Q2", "Prototype", "MVP ready"),
        ("Q3", "Launch", "Go live"),
    ])
    builder.add_quote_slide("The best way to predict the future is to create it.", "Peter Drucker")
    builder.add_kpi_dashboard("KPIs", [
        ("Revenue", "$1.2M", "+15%"),
        ("Users", "42K", "+8%"),
        ("NPS", "72", "+5"),
    ])
    path = builder.save("/tmp/demo_deck.pptx")
    print(f"Saved to {path}")
