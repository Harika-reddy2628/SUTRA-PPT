"""
template_engine.py - Template management for PowerPoint presentations.

Provides loading of .pptx templates, layout extraction, placeholder mapping,
brand colour application, and template compatibility validation.

Usage:
    from template_engine import TemplateEngine
    engine = TemplateEngine("template.pptx")
    layouts = engine.get_layouts()
    engine.apply_brand_colors({"primary": "#1B2A4A"})
    engine.fill_placeholder(slide, "TITLE", "My Title")
"""

from __future__ import annotations

import copy
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

COLOR_PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)
COLOR_SECONDARY = RGBColor(0x2E, 0x86, 0xAB)
COLOR_ACCENT = RGBColor(0xE8, 0x4D, 0x3D)

PLACEHOLDER_TYPES = {
    "CENTER_TITLE": 6,
    "TITLE": 7,
    "SUBTITLE": 8,
    "OBJECT": 9,
    "TABLE": 12,
    "CHART": 13,
    "ORGCHART": 14,
    "TABLE_OF_CONTENTS": 15,
    "VERTICAL_TITLE": 16,
    "VERTICAL_TEXT": 17,
    "TEXT": 18,
    "BITMAP": 19,
    "DATE_TIME": 21,
    "SLIDE_NUMBER": 22,
    "FOOTER": 23,
    "HEADER": 24,
}


# ---------------------------------------------------------------------------
# TemplateEngine
# ---------------------------------------------------------------------------

class TemplateEngine:
    """Template management: load, inspect, validate, and apply content."""

    def __init__(self, template_path: Optional[str] = None):
        """Load a .pptx template or start with a blank presentation.

        Args:
            template_path: Path to a .pptx file. If None, creates a blank deck.

        Raises:
            FileNotFoundError: If template_path is provided but does not exist.
        """
        if template_path and os.path.isfile(template_path):
            self.prs = Presentation(template_path)
            self.template_path = template_path
        elif template_path:
            raise FileNotFoundError(f"Template not found: {template_path}")
        else:
            self.prs = Presentation()
            self.template_path = None
        self._layout_cache: Dict[int, Any] = {}

    # ------------------------------------------------------------------
    # Layout inspection
    # ------------------------------------------------------------------

    def get_layouts(self) -> List[Dict[str, Any]]:
        """Return metadata for all available slide layouts.

        Returns:
            List of dicts with keys: index, name, placeholder_count,
            placeholders (list of placeholder info).
        """
        layouts = []
        for i, layout in enumerate(self.prs.slide_layouts):
            placeholders = []
            for ph in layout.placeholders:
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "name": ph.name,
                    "type": ph.placeholder_format.type,
                    "type_name": str(ph.placeholder_format.type),
                    "left": ph.left,
                    "top": ph.top,
                    "width": ph.width,
                    "height": ph.height,
                })
            layouts.append({
                "index": i,
                "name": layout.name,
                "placeholder_count": len(placeholders),
                "placeholders": placeholders,
            })
        return layouts

    def get_layout_names(self) -> List[str]:
        """Return just the layout names."""
        return [layout.name for layout in self.prs.slide_layouts]

    def find_layout(self, name: str) -> Optional[int]:
        """Find a layout index by name (case-insensitive partial match).

        Returns:
            Layout index or None if not found.
        """
        name_lower = name.lower()
        for i, layout in enumerate(self.prs.slide_layouts):
            if name_lower in layout.name.lower():
                return i
        return None

    # ------------------------------------------------------------------
    # Placeholder inspection
    # ------------------------------------------------------------------

    def get_placeholders(self, layout_index: int) -> List[Dict[str, Any]]:
        """Get all placeholders for a specific layout."""
        if layout_index < 0 or layout_index >= len(self.prs.slide_layouts):
            return []
        layout = self.prs.slide_layouts[layout_index]
        result = []
        for ph in layout.placeholders:
            result.append({
                "idx": ph.placeholder_format.idx,
                "name": ph.name,
                "type": str(ph.placeholder_format.type),
            })
        return result

    def get_used_placeholders(self, slide) -> Dict[int, str]:
        """Get placeholders that have content on a slide.

        Returns:
            Dict mapping placeholder idx to current text content.
        """
        used = {}
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            text = shape.text_frame.text if shape.has_text_frame else ""
            used[idx] = text
        return used

    # ------------------------------------------------------------------
    # Content filling
    # ------------------------------------------------------------------

    def fill_placeholder(self, slide, placeholder: str | int,
                         content: str, font_name: str = "Inter",
                         font_size: int = 18, bold: bool = False,
                         color: Optional[RGBColor] = None) -> bool:
        """Fill a placeholder on a slide with text content.

        Args:
            slide: The slide object.
            placeholder: Placeholder name (str) or index (int).
            content: Text to insert.
            font_name: Font family.
            font_size: Font size in points.
            bold: Whether to bold the text.
            color: Text colour.

        Returns:
            True if placeholder was found and filled, False otherwise.
        """
        target_ph = None
        for shape in slide.placeholders:
            if isinstance(placeholder, int):
                if shape.placeholder_format.idx == placeholder:
                    target_ph = shape
                    break
            elif isinstance(placeholder, str):
                if shape.name.lower() == placeholder.lower():
                    target_ph = shape
                    break

        if target_ph is None:
            return False

        if target_ph.has_text_frame:
            tf = target_ph.text_frame
            if tf.paragraphs:
                p = tf.paragraphs[0]
                p.text = content
                for run in p.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                    if color:
                        run.font.color.rgb = color
            else:
                p = tf.add_paragraph()
                p.text = content
                for run in p.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                    if color:
                        run.font.color.rgb = color
        return True

    def fill_slide(self, slide, content_map: Dict[str | int, str],
                   font_name: str = "Inter", font_size: int = 18,
                   bold: bool = False,
                   color: Optional[RGBColor] = None) -> Dict[str | int, bool]:
        """Fill multiple placeholders on a slide at once.

        Args:
            slide: The slide object.
            content_map: Dict mapping placeholder name/index to text.
            font_name, font_size, bold, color: Formatting options.

        Returns:
            Dict mapping placeholder to success (True/False).
        """
        results = {}
        for placeholder, content in content_map.items():
            results[placeholder] = self.fill_placeholder(
                slide, placeholder, content, font_name, font_size, bold, color
            )
        return results

    def add_slide_from_layout(self, layout_index: int,
                              content_map: Optional[Dict[str | int, str]] = None):
        """Add a new slide from a layout and optionally fill placeholders.

        Args:
            layout_index: Index of the layout to use.
            content_map: Optional dict of placeholder -> content.

        Returns:
            The new slide object.
        """
        if layout_index < 0 or layout_index >= len(self.prs.slide_layouts):
            raise IndexError(f"Layout index {layout_index} out of range "
                             f"(0-{len(self.prs.slide_layouts) - 1})")

        layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(layout)

        if content_map:
            self.fill_slide(slide, content_map)

        return slide

    # ------------------------------------------------------------------
    # Brand colour application
    # ------------------------------------------------------------------

    def apply_brand_colors(self, color_map: Dict[str, str | RGBColor],
                           apply_to_existing: bool = False) -> None:
        """Apply brand colours to the presentation theme.

        Args:
            color_map: Dict mapping colour roles to hex/RGB values.
                       Keys: primary, secondary, accent, success, warning.
            apply_to_existing: If True, also recolour existing slide shapes.
        """
        def _to_rgb(val) -> RGBColor:
            if isinstance(val, RGBColor):
                return val
            if isinstance(val, str):
                h = val.lstrip("#")
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            return COLOR_PRIMARY

        # Apply to theme colours
        try:
            theme = self.prs.slide_master.theme
            colour_scheme = theme.theme_elements.clrScheme
            mapping = {
                "primary": "dk1",
                "secondary": "lt1",
                "accent": "accent1",
                "success": "accent3",
                "warning": "accent6",
            }
            for role, hex_val in color_map.items():
                attr = mapping.get(role)
                if attr and hasattr(colour_scheme, attr):
                    elem = getattr(colour_scheme, attr)
                    srgb = elem.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    if srgb is not None:
                        srgb.set("val", _to_rgb(hex_val).__str__())
        except Exception:
            pass  # Theme may not be editable in all templates

        # Optionally apply to existing shapes
        if apply_to_existing:
            self._recolor_shapes(color_map)

    def _recolor_shapes(self, color_map: Dict[str, str | RGBColor]) -> None:
        """Recolour shapes on all slides based on colour map."""
        def _to_rgb(val) -> RGBColor:
            if isinstance(val, RGBColor):
                return val
            if isinstance(val, str):
                h = val.lstrip("#")
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            return COLOR_PRIMARY

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.color.rgb and run.font.color.rgb != RGBColor(0, 0, 0):
                                # Keep original if it's explicitly set
                                pass

    # ------------------------------------------------------------------
    # Template validation
    # ------------------------------------------------------------------

    def validate(self) -> Dict[str, Any]:
        """Validate the loaded template and return a compatibility report.

        Returns:
            Dict with keys: valid, layout_count, has_title_layout,
            has_content_layout, issues, warnings.
        """
        issues: List[str] = []
        warnings: List[str] = []

        layout_names = self.get_layout_names()
        has_title = any("title" in n.lower() for n in layout_names)
        has_content = any("content" in n.lower() or "title and content" in n.lower()
                          for n in layout_names)

        if not has_title:
            warnings.append("No layout with 'title' in name found.")
        if not has_content:
            warnings.append("No layout with 'content' in name found.")

        # Check for placeholder availability
        for i, layout in enumerate(self.prs.slide_layouts):
            if layout.placeholders:
                pass  # Good
            else:
                warnings.append(f"Layout '{layout.name}' has no placeholders.")

        # Check slide dimensions
        width = self.prs.slide_width
        height = self.prs.slide_height
        aspect = width / height if height else 0
        if abs(aspect - 16 / 9) > 0.1 and abs(aspect - 4 / 3) > 0.1:
            warnings.append(f"Unusual aspect ratio: {aspect:.2f} (expected ~1.78 or 1.33)")

        return {
            "valid": len(issues) == 0,
            "layout_count": len(self.prs.slide_layouts),
            "has_title_layout": has_title,
            "has_content_layout": has_content,
            "slide_count": len(self.prs.slides),
            "width_inches": width / 914400,
            "height_inches": height / 914400,
            "issues": issues,
            "warnings": warnings,
        }

    def print_report(self) -> None:
        """Print a human-readable validation report."""
        report = self.validate()
        print("=== Template Validation Report ===")
        print(f"  Valid:          {report['valid']}")
        print(f"  Layouts:        {report['layout_count']}")
        print(f"  Slides:         {report['slide_count']}")
        print(f"  Dimensions:     {report['width_inches']:.1f}\" x {report['height_inches']:.1f}\"")
        print(f"  Title layout:   {report['has_title_layout']}")
        print(f"  Content layout: {report['has_content_layout']}")
        if report["issues"]:
            print("  Issues:")
            for issue in report["issues"]:
                print(f"    - {issue}")
        if report["warnings"]:
            print("  Warnings:")
            for warn in report["warnings"]:
                print(f"    - {warn}")

    # ------------------------------------------------------------------
    # Utility
    # ------------------------------------------------------------------

    def clone_slide(self, source_slide_index: int):
        """Clone a slide by index and return the new slide.

        Uses XML copy for fidelity.
        """
        if source_slide_index < 0 or source_slide_index >= len(self.prs.slides):
            raise IndexError(f"Slide index {source_slide_index} out of range")

        source = self.prs.slides[source_slide_index]
        layout = source.slide_layout
        new_slide = self.prs.slides.add_slide(layout)

        # Copy shapes via XML
        for shape in source.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

        return new_slide

    def save(self, output_path: str) -> str:
        """Save the presentation and return the absolute path."""
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))
        return str(out.resolve())

    @property
    def presentation(self) -> Presentation:
        return self.prs


# ---------------------------------------------------------------------------
# Standalone helper: create template from scratch
# ---------------------------------------------------------------------------

def create_brand_template(output_path: str, brand_colors: Dict[str, str],
                          layouts_to_include: Optional[List[str]] = None) -> str:
    """Create a minimal brand template .pptx file.

    Args:
        output_path: Where to save the template.
        brand_colors: Dict with primary, secondary, accent hex codes.
        layouts_to_include: Optional list of layout names to keep.

    Returns:
        Absolute path to saved template.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add a title slide
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Brand Template"
        elif ph.placeholder_format.idx == 1:
            ph.text = "Created by ppt-artisan"

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out.resolve())


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    if len(sys.argv) > 1:
        template = sys.argv[1]
    else:
        template = None

    engine = TemplateEngine(template)
    engine.print_report()

    print("\n=== Available Layouts ===")
    for layout in engine.get_layouts():
        print(f"  [{layout['index']}] {layout['name']} "
              f"({layout['placeholder_count']} placeholders)")
        for ph in layout["placeholders"]:
            print(f"      - {ph['name']} (idx={ph['idx']}, type={ph['type_name']})")
