"""
chart_factory.py - Chart creation utilities for PowerPoint.

Provides bar, line, pie, and combo chart creation with consistent styling.
All charts follow clean design: no gridlines by default, clean labels,
accent colour highlighting.

Usage:
    from chart_factory import ChartFactory
    factory = ChartFactory()
    chart = factory.create_bar_chart(data, labels, colors, title)
"""

from __future__ import annotations

from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LEGEND_POSITION,
    XL_LABEL_POSITION,
    XL_MARK_STYLE,
)
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Design constants (mirror deck_builder for consistency)
# ---------------------------------------------------------------------------

COLOR_PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)
COLOR_SECONDARY = RGBColor(0x2E, 0x86, 0xAB)
COLOR_ACCENT = RGBColor(0xE8, 0x4D, 0x3D)
COLOR_SUCCESS = RGBColor(0x27, 0xAE, 0x60)
COLOR_WARNING = RGBColor(0xF3, 0x9C, 0x12)
COLOR_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
COLOR_DARK = RGBColor(0x2C, 0x3E, 0x50)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DEFAULT_CHART_COLORS = [COLOR_PRIMARY, COLOR_SECONDARY, COLOR_ACCENT,
                        COLOR_SUCCESS, COLOR_WARNING]

FONT_HEADING = "Montserrat"
FONT_BODY = "Inter"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _ensure_rgb(color) -> RGBColor:
    """Ensure the value is an RGBColor."""
    if isinstance(color, RGBColor):
        return color
    if isinstance(color, str):
        h = color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return COLOR_PRIMARY


def _style_axis(axis, font_name: str = FONT_BODY, font_size: int = 10,
                color: RGBColor = COLOR_DARK) -> None:
    """Style a chart axis (font, tick labels)."""
    try:
        axis.tick_labels.font.name = font_name
        axis.tick_labels.font.size = Pt(font_size)
        axis.tick_labels.font.color.rgb = color
    except Exception:
        pass


def _remove_gridlines(chart, axis_index: int = 0) -> None:
    """Remove major gridlines from the value axis."""
    try:
        value_axis = chart.value_axis
        value_axis.major_gridlines.format.line.fill.background()
    except Exception:
        pass


def _set_no_line(axis) -> None:
    """Set axis line to no-fill."""
    try:
        axis.format.line.fill.background()
    except Exception:
        pass


# ---------------------------------------------------------------------------
# ChartFactory
# ---------------------------------------------------------------------------

class ChartFactory:
    """Factory for creating consistently styled PowerPoint charts."""

    def __init__(self, colors: Optional[List[RGBColor]] = None,
                 heading_font: str = FONT_HEADING,
                 body_font: str = FONT_BODY):
        """Initialise the factory with optional colour overrides."""
        self.colors = colors or DEFAULT_CHART_COLORS
        self.heading_font = heading_font
        self.body_font = body_font

    def _get_color(self, index: int) -> RGBColor:
        """Get colour at index, cycling if needed."""
        return self.colors[index % len(self.colors)]

    # ------------------------------------------------------------------
    # Bar chart
    # ------------------------------------------------------------------

    def create_bar_chart(self, data: List[List[float]], labels: List[str],
                         colors: Optional[List[RGBColor]] = None,
                         title: str = "",
                         series_names: Optional[List[str]] = None,
                         stacked: bool = False) -> Any:
        """Create a clustered or stacked bar chart.

        Args:
            data: List of series, each series is a list of values.
            labels: Category labels.
            colors: Optional per-series colours.
            title: Chart title text.
            series_names: Optional series names.
            stacked: If True, use stacked column chart.

        Returns:
            pptx.chart.Chart object.
        """
        chart_data = CategoryChartData()
        chart_data.categories = labels

        if series_names is None:
            series_names = [f"Series {i + 1}" for i in range(len(data))]
        for name, values in zip(series_names, data):
            chart_data.add_series(name, values)

        chart_type = (XL_CHART_TYPE.COLUMN_STACKED if stacked
                      else XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Return the chart_data and type for caller to add to slide
        return {
            "chart_data": chart_data,
            "chart_type": chart_type,
            "title": title,
            "colors": colors or self.colors,
            "series_count": len(data),
        }

    # ------------------------------------------------------------------
    # Line chart
    # ------------------------------------------------------------------

    def create_line_chart(self, data: List[List[float]], labels: List[str],
                          colors: Optional[List[RGBColor]] = None,
                          title: str = "",
                          series_names: Optional[List[str]] = None,
                          smooth: bool = True,
                          markers: bool = True) -> Any:
        """Create a line chart with optional smoothing and markers.

        Returns:
            Dict with chart_data, chart_type, title, colours.
        """
        chart_data = CategoryChartData()
        chart_data.categories = labels

        if series_names is None:
            series_names = [f"Series {i + 1}" for i in range(len(data))]
        for name, values in zip(series_names, data):
            chart_data.add_series(name, values)

        if smooth:
            chart_type = (XL_CHART_TYPE.LINE_MARKERS if markers
                          else XL_CHART_TYPE.LINE)
        else:
            chart_type = (XL_CHART_TYPE.LINE_MARKERS if markers
                          else XL_CHART_TYPE.LINE)

        return {
            "chart_data": chart_data,
            "chart_type": chart_type,
            "title": title,
            "colors": colors or self.colors,
            "series_count": len(data),
        }

    # ------------------------------------------------------------------
    # Pie chart
    # ------------------------------------------------------------------

    def create_pie_chart(self, data: List[float], labels: List[str],
                         colors: Optional[List[RGBColor]] = None,
                         title: str = "",
                         show_percentage: bool = True,
                         show_legend: bool = True) -> Any:
        """Create a pie chart.

        Returns:
            Dict with chart_data, chart_type, title, colours.
        """
        chart_data = ChartData()
        chart_data.categories = labels
        chart_data.add_series("Data", data)

        return {
            "chart_data": chart_data,
            "chart_type": XL_CHART_TYPE.PIE,
            "title": title,
            "colors": colors or self.colors,
            "series_count": 1,
            "show_percentage": show_percentage,
            "show_legend": show_legend,
        }

    # ------------------------------------------------------------------
    # Combo chart (column + line)
    # ------------------------------------------------------------------

    def create_combo_chart(self, data: List[List[float]], labels: List[str],
                           title: str = "",
                           series_names: Optional[List[str]] = None,
                           line_series_indices: Optional[List[int]] = None) -> Any:
        """Create a combo chart (column primary + line secondary).

        Returns:
            Dict with chart_data, chart_type, title.
        """
        chart_data = CategoryChartData()
        chart_data.categories = labels

        if series_names is None:
            series_names = [f"Series {i + 1}" for i in range(len(data))]
        for name, values in zip(series_names, data):
            chart_data.add_series(name, values)

        return {
            "chart_data": chart_data,
            "chart_type": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "title": title,
            "colors": self.colors,
            "series_count": len(data),
            "line_series_indices": line_series_indices or [1],
            "is_combo": True,
        }

    # ------------------------------------------------------------------
    # Apply chart to slide
    # ------------------------------------------------------------------

    def add_chart_to_slide(self, slide, chart_spec: Dict[str, Any],
                           left: Inches = Inches(1),
                           top: Inches = Inches(1.5),
                           width: Inches = Inches(11),
                           height: Inches = Inches(5)) -> Any:
        """Add a chart specification to a slide and style it.

        Args:
            slide: pptx slide object.
            chart_spec: Dict returned by create_*_chart methods.
            left, top, width, height: Position and size.

        Returns:
            The chart object.
        """
        chart_data = chart_spec["chart_data"]
        chart_type = chart_spec["chart_type"]
        colors = chart_spec.get("colors", self.colors)

        chart_frame = slide.shapes.add_chart(
            chart_type, left, top, width, height, chart_data
        )
        chart = chart_frame.chart

        # Style the chart
        self.style_chart(chart, colors, chart_spec)

        return chart

    # ------------------------------------------------------------------
    # Style chart
    # ------------------------------------------------------------------

    def style_chart(self, chart, colors: Optional[List[RGBColor]] = None,
                    spec: Optional[Dict[str, Any]] = None) -> None:
        """Apply consistent styling to a chart.

        - No gridlines
        - Accent colour per series
        - Clean axis labels
        - Legend positioning
        """
        colors = colors or self.colors
        spec = spec or {}

        # Remove gridlines
        _remove_gridlines(chart)

        # Style axes
        try:
            _style_axis(chart.value_axis, self.body_font, 10, COLOR_DARK)
            _set_no_line(chart.value_axis)
        except Exception:
            pass

        try:
            _style_axis(chart.category_axis, self.body_font, 10, COLOR_DARK)
            _set_no_line(chart.category_axis)
        except Exception:
            pass

        # Colour each series
        for i, series in enumerate(chart.series):
            color = colors[i % len(colors)]
            try:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
            except Exception:
                pass

            # Data labels
            try:
                series.has_data_labels = True
                data_labels = series.data_labels
                data_labels.font.name = self.body_font
                data_labels.font.size = Pt(9)
                data_labels.font.color.rgb = COLOR_DARK
                data_labels.number_format = "#,##0"
            except Exception:
                pass

        # Legend
        try:
            chart.has_legend = spec.get("series_count", 1) > 1
            if chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.name = self.body_font
                chart.legend.font.size = Pt(10)
                chart.legend.font.color.rgb = COLOR_DARK
        except Exception:
            pass

        # Pie-specific: data labels with percentages
        if spec.get("show_percentage") and chart.chart_type == XL_CHART_TYPE.PIE:
            try:
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.font.name = self.body_font
                data_labels.font.size = Pt(10)
                data_labels.number_format = "0.0%"
                data_labels.show_percentage = True
                data_labels.show_value = False
                data_labels.show_category_name = False
            except Exception:
                pass

        # Title styling
        try:
            chart.has_title = bool(spec.get("title"))
            if chart.has_title:
                chart.chart_title.text_frame.paragraphs[0].text = spec["title"]
                chart.chart_title.text_frame.paragraphs[0].font.name = self.heading_font
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
                chart.chart_title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
                chart.chart_title.text_frame.paragraphs[0].font.bold = True
        except Exception:
            pass

    # ------------------------------------------------------------------
    # Quick helpers for adding charts directly to slides
    # ------------------------------------------------------------------

    def bar_to_slide(self, slide, data: List[List[float]], labels: List[str],
                     title: str = "", colors: Optional[List[RGBColor]] = None,
                     series_names: Optional[List[str]] = None,
                     left: Inches = Inches(1), top: Inches = Inches(1.5),
                     width: Inches = Inches(11), height: Inches = Inches(5)) -> Any:
        """Create and add a bar chart in one call."""
        spec = self.create_bar_chart(data, labels, colors, title, series_names)
        return self.add_chart_to_slide(slide, spec, left, top, width, height)

    def line_to_slide(self, slide, data: List[List[float]], labels: List[str],
                      title: str = "", colors: Optional[List[RGBColor]] = None,
                      series_names: Optional[List[str]] = None,
                      left: Inches = Inches(1), top: Inches = Inches(1.5),
                      width: Inches = Inches(11), height: Inches = Inches(5)) -> Any:
        """Create and add a line chart in one call."""
        spec = self.create_line_chart(data, labels, colors, title, series_names)
        return self.add_chart_to_slide(slide, spec, left, top, width, height)

    def pie_to_slide(self, slide, data: List[float], labels: List[str],
                     title: str = "", colors: Optional[List[RGBColor]] = None,
                     left: Inches = Inches(1), top: Inches = Inches(1.5),
                     width: Inches = Inches(11), height: Inches = Inches(5)) -> Any:
        """Create and add a pie chart in one call."""
        spec = self.create_pie_chart(data, labels, colors, title)
        return self.add_chart_to_slide(slide, spec, left, top, width, height)


# ---------------------------------------------------------------------------
# Standalone helper: style an existing chart
# ---------------------------------------------------------------------------

def style_chart(chart, theme: Optional[Dict[str, Any]] = None) -> None:
    """Apply default styling to an existing chart object.

    This is a convenience function for styling charts created outside
    the ChartFactory (e.g., directly via python-pptx).
    """
    factory = ChartFactory()
    colors = None
    if theme:
        colors = [theme.get("primary", COLOR_PRIMARY),
                  theme.get("secondary", COLOR_SECONDARY),
                  theme.get("accent", COLOR_ACCENT)]
    factory.style_chart(chart, colors)


# ---------------------------------------------------------------------------
# CLI entry point for testing
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    from pptx import Presentation

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    factory = ChartFactory()

    # Test bar chart
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    factory.bar_to_slide(slide1,
                         [[20, 35, 30, 45, 50], [15, 25, 40, 35, 45]],
                         ["Q1", "Q2", "Q3", "Q4", "Q5"],
                         title="Revenue by Quarter",
                         series_names=["2024", "2025"])

    # Test line chart
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    factory.line_to_slide(slide2,
                          [[10, 22, 30, 40, 55], [5, 15, 25, 38, 50]],
                          ["Jan", "Feb", "Mar", "Apr", "May"],
                          title="Growth Trend",
                          series_names=["Users", "Revenue"])

    # Test pie chart
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    factory.pie_to_slide(slide3,
                         [40, 30, 20, 10],
                         ["Product A", "Product B", "Product C", "Product D"],
                         title="Market Share")

    prs.save("/tmp/test_charts.pptx")
    print("Charts saved to /tmp/test_charts.pptx")
